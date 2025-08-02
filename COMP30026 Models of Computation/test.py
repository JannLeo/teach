from pptx import Presentation
from pptx.util import Inches
import os

# 设置图片文件夹路径
image_folder = r"D:\创业\teacher\teach\COMP30026 Models of Computation\reademe.assets"

# 创建新的 PowerPoint 演示文稿
prs = Presentation()
slide_width = prs.slide_width
slide_height = prs.slide_height

# 遍历所有图片文件
for filename in sorted(os.listdir(image_folder)):
    if filename.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif")):
        # 添加空白幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 插入图片（铺满整个幻灯片）
        img_path = os.path.join(image_folder, filename)
        slide.shapes.add_picture(img_path, 0, 0, width=slide_width, height=slide_height)

# 保存为PPT文件
output_path = os.path.join(image_folder, "InsertedImages.pptx")
prs.save(output_path)
print(f"PPT 已成功保存到: {output_path}")
